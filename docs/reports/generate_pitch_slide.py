#!/usr/bin/env python3
"""One-slide Young Builders pitch deck for TradePulse — template layout preserved."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "TradePulse_Young_Builders_Pitch_One_Slide.pptx"
OUT_ALT = Path(__file__).resolve().parent / "TradePulse_Young_Builders_Pitch_One_Slide_v2.pptx"

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xC4, 0x5C, 0x26)
NAVY = RGBColor(0x0B, 0x1F, 0x33)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
LINE_PT = Pt(1.0)


def set_run(run, text: str, *, size: Pt, bold: bool = False, italic: bool = False, color=BLACK):
    run.text = text
    run.font.name = "Arial"
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_box(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLACK
    shape.line.width = LINE_PT
    return shape


def set_shape_text(shape, paragraphs: list[tuple[str, dict]], *, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    try:
        from pptx.enum.text import MSO_AUTO_SIZE

        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf._txBody.bodyPr.set("anchor", "t")
    except Exception:
        pass

    first = True
    for text, opts in paragraphs:
        p = shape.text_frame.paragraphs[0] if first else shape.text_frame.add_paragraph()
        first = False
        p.alignment = opts.get("align", align)
        p.space_before = Pt(opts.get("space_before", 0))
        p.space_after = Pt(opts.get("space_after", 2))
        p.level = opts.get("level", 0)
        run = p.add_run()
        set_run(
            run,
            text,
            size=Pt(opts.get("size", 10)),
            bold=opts.get("bold", False),
            italic=opts.get("italic", False),
            color=opts.get("color", BLACK),
        )


def add_textbox(slide, left, top, width, height, paragraphs, *, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_shape_text(box, paragraphs, align=align)
    return box


def checkbox(marked: bool) -> str:
    return "■" if marked else "□"


def flow_node(slide, left, top, width, height, label: str, *, fill=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = NAVY
    shape.line.width = Pt(1.0)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    try:
        tf._txBody.bodyPr.set("anchor", "ctr")
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, label, size=Pt(8), bold=True, color=NAVY)
    return shape


def flow_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    return shape


def add_solution_flowchart(slide, col_left, col_top, col_w):
    """Horizontal flow inside the Solution column (below short copy)."""
    pad = Inches(0.12)
    y = col_top + Inches(1.05)
    node_h = Inches(0.58)
    arrow_w = Inches(0.22)
    arrow_h = Inches(0.16)
    inner_w = col_w - 2 * pad
    # 4 nodes + 3 arrows
    node_w = (inner_w - 3 * arrow_w) / 4
    x = col_left + pad
    labels = ["Docs in", "Double-\ncheck", "Flag\ngaps", "Officer\ndecides"]
    for i, label in enumerate(labels):
        flow_node(slide, x, y, node_w, node_h, label)
        x += node_w
        if i < len(labels) - 1:
            flow_arrow(
                slide,
                x,
                y + (node_h - arrow_h) / 2,
                arrow_w,
                arrow_h,
            )
            x += arrow_w

    # Caption under flow
    add_textbox(
        slide,
        col_left + pad,
        y + node_h + Inches(0.06),
        col_w - 2 * pad,
        Inches(0.35),
        [
            (
                "Humans stay in charge — TradePulse prepares the file",
                {"size": 7.5, "italic": True, "align": PP_ALIGN.CENTER, "color": ORANGE},
            )
        ],
        align=PP_ALIGN.CENTER,
    )


def add_evidence_chips(slide, col_left, col_top, col_w):
    """Three small proof chips in Validation column."""
    pad = Inches(0.12)
    y = col_top + Inches(1.35)
    chip_h = Inches(0.55)
    gap = Inches(0.08)
    chip_w = (col_w - 2 * pad - 2 * gap) / 3
    chips = ["Live\ndemo", "Sample\ncases", "Handoff\npack"]
    x = col_left + pad
    for label in chips:
        flow_node(slide, x, y, chip_w, chip_h, label, fill=WHITE)
        x += chip_w + gap


def add_gift_flowchart(slide, box_left, box_top, box_w):
    """Why GIFT IFIH — visual path from residency to pilot."""
    pad = Inches(0.18)
    y = box_top + Inches(0.38)
    node_h = Inches(0.42)
    arrow_w = Inches(0.28)
    arrow_h = Inches(0.14)
    inner_w = box_w - 2 * pad
    labels = [
        "Mentors &\nregulators",
        "GIFT trade\ndesk intros",
        "Sandbox /\npilot path",
        "Bank Ops\npilot",
    ]
    node_w = (inner_w - 3 * arrow_w) / 4
    x = box_left + pad
    for i, label in enumerate(labels):
        flow_node(slide, x, y, node_w, node_h, label, fill=WHITE)
        x += node_w
        if i < len(labels) - 1:
            flow_arrow(slide, x, y + (node_h - arrow_h) / 2, arrow_w, arrow_h)
            x += arrow_w

    add_textbox(
        slide,
        box_left + pad,
        y + node_h + Inches(0.04),
        box_w - 2 * pad,
        Inches(0.28),
        [
            (
                "Young Builders → residency + IFSCA exposure → convert tonight’s prototype into a supervised bank pilot",
                {"size": 7.5, "italic": True, "align": PP_ALIGN.CENTER, "color": ORANGE},
            )
        ],
        align=PP_ALIGN.CENTER,
    )



def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    L = Inches(0.35)
    R = Inches(0.35)
    usable_w = prs.slide_width - L - R
    gap = Inches(0.12)

    # --- Header (preserved) ---
    add_textbox(
        slide,
        L,
        Inches(0.12),
        usable_w,
        Inches(0.38),
        [
            (
                "Young Builders Program Hackathon | Pitch Template (One Slide)",
                {"size": 16, "bold": True, "align": PP_ALIGN.CENTER},
            )
        ],
        align=PP_ALIGN.CENTER,
    )

    # --- Pitch Snapshot ---
    y = Inches(0.52)
    h_snap = Inches(0.95)
    snap = add_box(slide, L, y, usable_w, h_snap)
    set_shape_text(
        snap,
        [
            ("Pitch Snapshot", {"size": 12, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 5}),
            (
                "• Team: TradePulse  |  Track 1 — Agentic AI  |  Focus: Faster, safer trade-document review for bank / GIFT officers",
                {"size": 9, "space_after": 3},
            ),
            (
                f"• Stage: {checkbox(False)} Idea   {checkbox(True)} Prototype   {checkbox(False)} MVP   {checkbox(False)} Pilot"
                f"     • Commitment: {checkbox(False)} Full-time   {checkbox(True)} Part-time",
                {"size": 9, "space_after": 0},
            ),
        ],
    )

    # --- Row: Problem | Solution | Validation ---
    y = y + h_snap + gap
    h_mid = Inches(2.55)
    col_w = (usable_w - 2 * gap) / 3

    problem = add_box(slide, L, y, col_w, h_mid)
    set_shape_text(
        problem,
        [
            ("1. Problem", {"size": 11, "bold": True, "space_after": 6}),
            ("• Who: Trade officers at banks & GIFT City IBUs", {"size": 9, "space_after": 4}),
            ("• Buyer: Head of Trade Finance Ops", {"size": 9, "space_after": 4}),
            ("• Today: Hand-checking invoice + shipping PDFs", {"size": 9, "space_after": 4}),
            ("• Trap: Look-alike names treated as “confirmed”", {"size": 9, "space_after": 4}),
            ("• Pain: Hours per file · weak audit trail · rising volume", {"size": 9, "space_after": 0}),
        ],
    )

    sol_left = L + col_w + gap
    solution = add_box(slide, sol_left, y, col_w, h_mid)
    set_shape_text(
        solution,
        [
            ("2. Solution", {"size": 11, "bold": True, "space_after": 5}),
            ("• TradePulse = review desk (not auto-approve)", {"size": 9, "space_after": 3}),
            ("• Missing data never becomes a fake “all clear”", {"size": 9, "space_after": 2}),
        ],
    )
    add_solution_flowchart(slide, sol_left, y, col_w)

    val_left = L + 2 * (col_w + gap)
    validation = add_box(slide, val_left, y, col_w, h_mid)
    set_shape_text(
        validation,
        [
            ("3. Validation", {"size": 11, "bold": True, "space_after": 5}),
            ("• Buyer locked: GIFT IBU / bank trade desk Ops head", {"size": 9, "space_after": 3}),
            ("• Banks review paperwork — not the physical cargo", {"size": 9, "space_after": 3}),
            ("• Proof: live AWS demo (Mumbai), demo data labelled", {"size": 9, "space_after": 2}),
        ],
    )
    add_evidence_chips(slide, val_left, y, col_w)

    # --- Row: Regulatory | Team-Market Fit ---
    y = y + h_mid + gap
    h_low = Inches(1.50)
    half_w = (usable_w - gap) / 2

    regulatory = add_box(slide, L, y, half_w, h_low)
    set_shape_text(
        regulatory,
        [
            ("4. Regulatory Requirements (Domestic/IFSC)", {"size": 11, "bold": True, "space_after": 4}),
            ("• Focus: IFSCA / GIFT City trade desks; bank trade ops under RBI", {"size": 8.5, "space_after": 2}),
            ("• Company ID: use official records when available — look-alike names are not proof", {"size": 8.5, "space_after": 2}),
            ("• We never claim: AI approved / cleared / sanctioned a trade", {"size": 8.5, "space_after": 2}),
            ("• Missing papers stay visible — officers stay in the loop on every material gap", {"size": 8.5, "space_after": 0}),
        ],
    )

    team = add_box(slide, L + half_w + gap, y, half_w, h_low)
    set_shape_text(
        team,
        [
            ("5. Team-Market Fit", {"size": 11, "bold": True, "space_after": 4}),
            ("• Abhishek — platform  ·  Ansh — product desk  ·  Atharva — UI  ·  Shivansh — QA", {"size": 8.5, "space_after": 2}),
            ("• Clear owners for every part — we stop when rules conflict, we don’t invent mid-demo", {"size": 8.5, "space_after": 2}),
            ("• Built for the Ops buyer, not a generic “AI chatbot for PDFs” demo", {"size": 8.5, "space_after": 2}),
            ("• Biggest win: live cloud review desk officers can open and click today", {"size": 8.5, "space_after": 0}),
        ],
    )

    # --- Why GIFT IFIH? ---
    y = y + h_low + gap
    h_gift = Inches(1.25)
    gift = add_box(slide, L, y, usable_w, h_gift)
    set_shape_text(
        gift,
        [
            ("Why GIFT IFIH?", {"size": 12, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 2}),
        ],
    )
    add_gift_flowchart(slide, L, y, usable_w)

    # --- Footer (preserved) ---
    add_textbox(
        slide,
        L,
        Inches(7.12),
        usable_w,
        Inches(0.32),
        [
            (
                "This is a guide for the content of your slide. Structure, formatting, and design are left to your discretion. Use evidence & metrics wherever possible. Avoid paragraphs.",
                {"size": 8, "italic": True, "color": ORANGE, "align": PP_ALIGN.LEFT},
            )
        ],
    )

    target = OUT
    try:
        with open(OUT, "a"):
            pass
    except PermissionError:
        target = OUT_ALT
    prs.save(target)
    print(target)


if __name__ == "__main__":
    main()
